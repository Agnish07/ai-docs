# backend/app/core/export_pptx.py
import io
import json
from pptx import Presentation
from pptx.util import Inches, Pt

def build_pptx_bytes(project, db=None):
    """
    Minimal pptx exporter using python-pptx.
    Creates one slide per item with title + content paragraph.
    Requires python-pptx: pip install python-pptx
    """
    prs = Presentation()
    # choose a simple blank layout (0 usually has title+subtitle; 1 often title+content)
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    items = getattr(project, "items", []) or []
    items = sorted(items, key=lambda it: (getattr(it, "order", 0) or 0))

    for it in items:
        slide = prs.slides.add_slide(layout)
        # set title if placeholder exists
        try:
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = getattr(it, "title", "") or ""
        except Exception:
            pass

        # get content
        content_md = ""
        try:
            raw = getattr(it, "content", "") or ""
            parsed = json.loads(raw) if isinstance(raw, str) and raw.strip().startswith("{") else None
            if parsed and parsed.get("sections"):
                content_md = parsed["sections"][0].get("content_md", "")
            else:
                content_md = raw or ""
        except Exception:
            content_md = getattr(it, "content", "") or ""

        # add text box (fallback if there is no body placeholder)
        body_placeholder = None
        for shape in slide.placeholders:
            phf = getattr(shape, "placeholder_format", None)
            if phf and phf.type is not None:
                body_placeholder = shape
                break

        if body_placeholder and body_placeholder.has_text_frame:
            tf = body_placeholder.text_frame
            tf.text = content_md or ""
        else:
            # fallback: make a textbox
            left = Inches(1)
            top = Inches(1.6)
            width = Inches(8)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = content_md or ""

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()
